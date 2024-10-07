import os
import pptx
from pptx import Presentation
from pptx.util import Inches
from model import GPT
from utils import add_content

class SlidesGenerator:
    def __init__(self) -> None:
        self.gpt = GPT('gpt-4o')


    def generate_raw_text(self, title, contents, slides_number, folder):
        cache_file = f'Cache/{folder}_raw.txt'
        if os.path.exists(cache_file):
            with open(cache_file, 'r', encoding='utf-8') as f:
                slides_text = f.read()
        else:
            paper = [add_content(section[0], section[1]) for section in contents]
            prompt = f""" Write a presentation text for a scientific paper. You should only provide the finished presentation.
Follow these guidelines:
- The text must not exceed 1000 characters.
- Use the original title of the paper as the presentation title.
- The presentation must include a table of contents that matches the number of slides/sections.
- Include a summary at the end.
- Do not insert any links or images.
- The presentation must contain exactly {slides_number} slides/sections.

Title of the paper: {title}
Content of the paper: {paper}

Example format: (Stick to this format exactly!)
```
Title: TITLE OF THE PRESENTATION

Slide: 1
Header: TABLE OF CONTENT
Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENT OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

Slide: 2
Header: TITLE OF SLIDE
Content: CONTENT OF THE SLIDE

Slide: 3
Header: TITLE OF SLIDE
Content: CONTENT OF THE SLIDE

...

Slide: X
Headers: SUMMARY
Content: CONTENT OF THE SUMMARY

Slide: END
```
"""
            result = self.gpt.infer(user_prompt = prompt)
            slides_text = result
        with open(cache_file, 'w', encoding='utf-8') as f:
            f.write(slides_text)
        return slides_text


    def generate_slide(self, slides_text, theme, file_name, folder):
        prs = Presentation(f"Designs/Design-{theme}.pptx")
        slide_count = 0
        header = ""
        content = ""
        lines = slides_text.splitlines()
        line_num = -1
        while line_num < len(lines) - 1:
            line_num += 1
            line = lines[line_num]
            if line.startswith('Title:'):
                header = line.replace('Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = header
                    body_shape = slide.shapes.placeholders[1]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" # reset content
                slide_count += 1 # add new slide
                continue
            elif line.startswith('Header:'):
                header = line.replace('Header:', '').strip()
                continue
            elif line.startswith('Content:'):
                content = line.replace('Content:', '').strip()
                line_num += 1

                while line_num < len(lines):
                    next_line = lines[line_num].strip()
                    if next_line.startswith('Figure:') or len(next_line) == 0:
                        break
                    content += '\n' + next_line
                    line_num += 1

                if next_line.startswith('Figure:'):
                    try:
                        figure_path = next_line.replace('Figure:', '').strip()
                        figure_path = os.path.join(folder, figure_path)
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        slide.shapes.add_picture(figure_path, pptx.util.Inches(2), pptx.util.Inches(2))
                    except:
                        pass
                continue

        prs.save(f'Result/{file_name}.pptx')
        file_path = f"Result/{file_name}.pptx"
        return str(f"{file_path}")

